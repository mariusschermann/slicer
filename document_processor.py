import os
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter
from docx2pdf import convert
import logging
from io import BytesIO

# Configure logging
logging.basicConfig(level=logging.INFO,
                   format='%(asctime)s - %(message)s',
                   datefmt='%Y-%m-%d %H:%M:%S')

# Global constants
SIZE_THRESHOLD = 15 * 1024 * 1024  # 15 MB in bytes
CHAR_THRESHOLD = 60000
WATCH_DIRECTORY = "/Users/schermannmarius/Test_Slicer"

def count_chars_docx(file_path):
    """Count characters in a DOCX file."""
    logging.info(f"Counting characters in DOCX: {file_path}")
    doc = Document(file_path)
    char_count = sum(len(paragraph.text) for paragraph in doc.paragraphs)
    logging.info(f"Character count for {file_path}: {char_count}")
    return char_count

def count_chars_pptx(file_path):
    """Count characters in a PPTX file."""
    logging.info(f"Counting characters in PPTX: {file_path}")
    prs = Presentation(file_path)
    char_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                char_count += len(shape.text)
    logging.info(f"Character count for {file_path}: {char_count}")
    return char_count

def count_chars_pdf(file_path):
    """Count characters in a PDF file."""
    logging.info(f"Counting characters in PDF: {file_path}")
    reader = PdfReader(file_path)
    char_count = 0
    for page in reader.pages:
        text = page.extract_text()
        char_count += len(text)
    logging.info(f"Character count for {file_path}: {char_count}")
    return char_count

def convert_to_pdf(file_path):
    """Convert DOCX/PPTX to PDF."""
    logging.info(f"Converting {file_path} to PDF")
    try:
        pdf_path = os.path.splitext(file_path)[0] + ".pdf"
        
        # Convert DOCX to PDF using docx2pdf
        if file_path.lower().endswith('.docx'):
            convert(file_path, pdf_path)
        # For PPTX files, we still need an alternative solution
        elif file_path.lower().endswith('.pptx'):
            logging.error("PPTX conversion not yet implemented")
            return None
            
        if os.path.exists(pdf_path):
            logging.info(f"Successfully converted {file_path} to PDF")
            return pdf_path
        else:
            logging.error(f"PDF conversion failed for {file_path}")
            return None
            
    except Exception as e:
        logging.error(f"Error converting file to PDF: {e}")
        return None

def slice_pdf(pdf_path):
    """Slice PDF into smaller parts based on size and character thresholds."""
    logging.info(f"Starting to slice PDF: {pdf_path}")
    reader = PdfReader(pdf_path)
    total_pages = len(reader.pages)
    current_size = 0
    current_chars = 0
    current_writer = PdfWriter()
    start_page_num = 1
    slices_created = False
    
    for i in range(total_pages):
        # Get page content and metrics before adding
        page = reader.pages[i]
        page_text = page.extract_text()
        page_chars = len(page_text)
        
        # Create a temporary PDF writer to measure the size
        temp_writer = PdfWriter()
        temp_writer.add_page(page)
        temp_buffer = BytesIO()
        temp_writer.write(temp_buffer)
        page_size = temp_buffer.tell()
        
        # Check if adding this page would exceed either limit
        if (current_chars + page_chars > CHAR_THRESHOLD or 
            current_size + page_size > SIZE_THRESHOLD):
            
            # If we have pages in the current writer, save them
            if current_writer.pages:
                output_path = f"{os.path.splitext(pdf_path)[0]}_{start_page_num}-{i}.pdf"
                with open(output_path, "wb") as output_file:
                    current_writer.write(output_file)
                logging.info(f"Created slice {output_path} (Pages: {start_page_num}-{i}, Size: {current_size/1024/1024:.2f}MB, Characters: {current_chars})")
                slices_created = True
                
                # Reset for next slice
                current_writer = PdfWriter()
                current_size = 0
                current_chars = 0
                start_page_num = i + 1
            
            # If a single page exceeds the limits, we need to handle it specially
            if page_chars > CHAR_THRESHOLD:
                logging.warning(f"Page {i+1} contains {page_chars} characters, which exceeds the limit of {CHAR_THRESHOLD}.")
                single_page_writer = PdfWriter()
                single_page_writer.add_page(page)
                output_path = f"{os.path.splitext(pdf_path)[0]}_{i+1}-{i+1}.pdf"
                with open(output_path, "wb") as output_file:
                    single_page_writer.write(output_file)
                logging.warning(f"Created oversized slice {output_path} (Single page {i+1}, Characters: {page_chars})")
                slices_created = True
                start_page_num = i + 2
                continue
        
        # Add page to current writer
        current_writer.add_page(page)
        current_size += page_size
        current_chars += page_chars
        
        # If this is the last page, save the current part
        if i == total_pages - 1 and current_writer.pages:
            output_path = f"{os.path.splitext(pdf_path)[0]}_{start_page_num}-{i+1}.pdf"
            with open(output_path, "wb") as output_file:
                current_writer.write(output_file)
            logging.info(f"Created final slice {output_path} (Pages: {start_page_num}-{i+1}, Size: {current_size/1024/1024:.2f}MB, Characters: {current_chars})")
            slices_created = True
    
    # Delete the original PDF if slices were created
    if slices_created:
        try:
            os.remove(pdf_path)
            logging.info(f"Deleted original PDF: {pdf_path}")
        except Exception as e:
            logging.error(f"Error deleting original PDF {pdf_path}: {e}")
    
    logging.info(f"Finished slicing PDF")

def process_file(file_path):
    """Process a single file."""
    logging.info(f"Processing file: {file_path}")
    _, ext = os.path.splitext(file_path)
    file_size = os.path.getsize(file_path)
    logging.info(f"File size: {file_size/1024/1024:.2f}MB")
    
    if ext.lower() in ['.docx', '.pptx']:
        char_count = 0
        if ext.lower() == '.docx':
            char_count = count_chars_docx(file_path)
        else:
            char_count = count_chars_pptx(file_path)
            
        logging.info(f"Checking thresholds - Size: {file_size/1024/1024:.2f}MB (threshold: {SIZE_THRESHOLD/1024/1024:.2f}MB), Characters: {char_count} (threshold: {CHAR_THRESHOLD})")
        
        if file_size > SIZE_THRESHOLD or char_count > CHAR_THRESHOLD:
            logging.info(f"File {file_path} exceeds thresholds, converting to PDF")
            pdf_path = convert_to_pdf(file_path)
            if pdf_path:
                slice_pdf(pdf_path)
        else:
            logging.info(f"File {file_path} does not exceed thresholds, skipping")
            
    elif ext.lower() == '.pdf':
        char_count = count_chars_pdf(file_path)
        logging.info(f"Checking thresholds - Size: {file_size/1024/1024:.2f}MB (threshold: {SIZE_THRESHOLD/1024/1024:.2f}MB), Characters: {char_count} (threshold: {CHAR_THRESHOLD})")
        if file_size > SIZE_THRESHOLD or char_count > CHAR_THRESHOLD:
            logging.info(f"Processing PDF {file_path} (Size: {file_size/1024/1024:.2f}MB, Characters: {char_count})")
            slice_pdf(file_path)
        else:
            logging.info(f"File {file_path} does not exceed thresholds, skipping")

def on_created(event):
    """Handle file creation events."""
    if not event.is_directory:
        logging.info(f"New file detected: {event.src_path}")
        process_file(event.src_path)

def main():
    """Main function to start the file monitoring."""
    if not os.path.exists(WATCH_DIRECTORY):
        print("Directory does not exist!")
        return

    logging.info(f"Starting document processor")
    logging.info(f"Watching directory: {WATCH_DIRECTORY}")
    logging.info(f"Size threshold: {SIZE_THRESHOLD/1024/1024:.2f}MB")
    logging.info(f"Character threshold: {CHAR_THRESHOLD}")
    
    # Process existing files
    for filename in os.listdir(WATCH_DIRECTORY):
        file_path = os.path.join(WATCH_DIRECTORY, filename)
        if os.path.isfile(file_path):
            process_file(file_path)

    # Set up the watchdog observer
    event_handler = FileSystemEventHandler()
    event_handler.on_created = on_created
    observer = Observer()
    observer.schedule(event_handler, WATCH_DIRECTORY, recursive=False)
    observer.start()

    logging.info("Press Ctrl+C to stop monitoring")

    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()

if __name__ == "__main__":
    main() 